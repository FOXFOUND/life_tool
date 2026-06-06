#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
自动为 PPTX 幻灯片添加备注。
用法：python add_notes_to_pptx.py <演示文稿路径> <备注文件路径> [输出路径]
备注文件每行格式：第X页：备注内容
"""

import sys
import re
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("错误：请先安装 python-pptx 库：pip install python-pptx")
    sys.exit(1)


def parse_notes_file(txt_path):
    """
    解析备注文本文件，返回备注列表（索引 0 对应第 1 页）。
    文件每行格式：第X页：备注内容（X 为页码）
    """
    notes = []
    line_pattern = re.compile(r'^第(\d+)页[：:](.*)$')

    with open(txt_path, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f, 1):
            line = line.strip()
            if not line:
                continue
            match = line_pattern.match(line)
            if not match:
                print(f"警告：第 {line_num} 行格式不正确，已跳过：{line[:50]}...")
                continue
            page_num = int(match.group(1))
            content = match.group(2).strip()
            # 确保列表长度足够，按页码存放（页码从 1 开始）
            while len(notes) < page_num:
                notes.append("")
            notes[page_num - 1] = content
    return notes


def add_notes_to_pptx(pptx_path, notes_list, output_path):
    """
    将备注列表写入 PPTX 的对应幻灯片备注页。
    """
    if not notes_list:
        print("警告：备注列表为空，无需处理。")
        return

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"演示文稿共有 {total_slides} 张幻灯片，备注文件提供了 {len(notes_list)} 条备注。")

    for idx, slide in enumerate(prs.slides):
        if idx >= len(notes_list):
            break
        note_text = notes_list[idx]
        if not note_text:
            continue

        # 获取备注页（如果不存在，notes_slide 属性会自动创建一个空的）
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = note_text
        print(f"已为第 {idx+1} 张幻灯片添加备注。")

    # 保存新文件
    prs.save(output_path)
    print(f"备注添加完成，已保存至：{output_path}")


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    pptx_path = sys.argv[1]
    notes_path = sys.argv[2]

    # 检查输入文件是否存在
    if not Path(pptx_path).is_file():
        print(f"错误：找不到演示文稿文件：{pptx_path}")
        sys.exit(1)
    if not Path(notes_path).is_file():
        print(f"错误：找不到备注文件：{notes_path}")
        sys.exit(1)

    # 确定输出路径
    if len(sys.argv) >= 4:
        output_path = sys.argv[3]
    else:
        base = Path(pptx_path).stem
        parent = Path(pptx_path).parent
        output_path = parent / f"{base}_with_notes.pptx"

    # 解析备注
    notes_list = parse_notes_file(notes_path)
    if not notes_list:
        print("错误：未从备注文件中解析到任何有效备注。")
        sys.exit(1)

    # 写入 PPTX
    add_notes_to_pptx(pptx_path, notes_list, output_path)


if __name__ == "__main__":
    main()