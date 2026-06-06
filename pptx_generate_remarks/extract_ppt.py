import sys
from pptx import Presentation

def extract_text(shape):
    """递归提取所有文本，兼容组合形状、表格、文本框"""
    text_content = []
    
    # 处理组合形状
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            text_content.extend(extract_text(sub_shape))
    
    # 普通文本
    if hasattr(shape, "text") and shape.text.strip():
        text_content.append(shape.text.strip())
    
    # 表格处理
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_data:
                text_content.append(" | ".join(row_data))
    
    return text_content

def extract_ppt_to_md(ppt_path, md_path="ppt_content.md"):
    prs = Presentation(ppt_path)
    
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("# PPT 文本内容提取\n\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            all_text = []
            for shape in slide.shapes:
                all_text.extend(extract_text(shape))
            
            # 去重、过滤空内容
            all_text = list(dict.fromkeys([t.strip() for t in all_text if t.strip()]))
            
            if not all_text:
                f.write(f"## 第 {slide_num} 页\n")
                f.write("无内容\n\n")
                continue
            
            # 标题 + 内容
            title = all_text[0]
            contents = all_text[1:]
            
            # Markdown 标准格式
            f.write(f"## 第 {slide_num} 页：{title}\n\n")
            for line in contents:
                f.write(f"- {line}\n")
            f.write("\n")
    
    return md_path

# ==================== 运行脚本 ====================
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python extract_ppt.py 隐形的镣铐.pptx")
        sys.exit(1)

    ppt_file = sys.argv[1]
    try:
        md_file = extract_ppt_to_md(ppt_file)
        print(f"✅ 提取完成！Markdown 文件已保存：{md_file}")
    except Exception as e:
        print(f"❌ 错误：{e}")